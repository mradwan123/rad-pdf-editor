"""Unit tests for core/ops/convert_from.py (PDF -> Word/PowerPoint/
Excel/HTML/JPG). All five operations here are pure-Python only (see
the module's docstring for why LibreOffice has no role in this
direction), so these tests are deterministic across machines - no
skipif needed."""

from __future__ import annotations

from pathlib import Path

import docx
import fitz
import openpyxl
import pytest
from pptx import Presentation

from core.errors import OperationError
from core.model.document import DocumentSession
from core.ops.convert_from import (
    PdfToDocxOperation,
    PdfToHtmlOperation,
    PdfToJpgOperation,
    PdfToPptxOperation,
    PdfToXlsxOperation,
)


def _session_with_pdf(tmp_path: Path, pages: list[tuple[str, tuple[float, float, float] | None]]) -> DocumentSession:
    """`pages`: list of (text, optional RGB fill for a marker rect)."""
    session_dir = tmp_path / "session"
    session_dir.mkdir(exist_ok=True)
    doc = fitz.open()
    for text, color in pages:
        page = doc.new_page(width=300, height=400)
        page.insert_text((50, 50), text)
        if color is not None:
            page.draw_rect(fitz.Rect(50, 100, 150, 150), color=color, fill=color)
    working = session_dir / "working.pdf"
    doc.save(working)
    doc.close()
    return DocumentSession(working_path=working, source_path=None)


# --- PdfToDocxOperation ---------------------------------------------------


def test_pdf_to_docx_produces_readable_paragraphs(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("Hello Word", None), ("Second page", None)])
    result = session.apply(PdfToDocxOperation())
    assert result.working_path.suffix == ".docx"

    document = docx.Document(str(result.working_path))
    text = "\n".join(p.text for p in document.paragraphs)
    assert "Hello Word" in text
    assert "Second page" in text


def test_pdf_to_docx_with_no_document_open_raises() -> None:
    doc = DocumentSession(working_path=None, source_path=None)
    with pytest.raises(OperationError):
        doc.apply(PdfToDocxOperation())


def test_pdf_to_docx_undo_restores_pdf(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("text", None)])
    result = session.apply(PdfToDocxOperation())
    restored = result.undo()
    assert restored.working_path.suffix == ".pdf"


# --- PdfToPptxOperation ---------------------------------------------------


def test_pdf_to_pptx_produces_one_slide_per_page(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("p1", None), ("p2", None), ("p3", None)])
    result = session.apply(PdfToPptxOperation(dpi=100))
    assert result.working_path.suffix == ".pptx"

    prs = Presentation(str(result.working_path))
    assert len(list(prs.slides)) == 3
    # each slide should have a full-slide picture shape
    for slide in prs.slides:
        assert len(slide.shapes) == 1


def test_pdf_to_pptx_rejects_too_low_dpi() -> None:
    with pytest.raises(OperationError):
        PdfToPptxOperation(dpi=10)


# --- PdfToXlsxOperation ----------------------------------------------------


def test_pdf_to_xlsx_extracts_table_cells(tmp_path: Path) -> None:
    session_dir = tmp_path / "session"
    session_dir.mkdir()
    doc = fitz.open()
    page = doc.new_page(width=300, height=400)
    # draw a simple 2x2 grid so pdfplumber detects a real table
    page.draw_line(fitz.Point(50, 50), fitz.Point(250, 50))
    page.draw_line(fitz.Point(50, 100), fitz.Point(250, 100))
    page.draw_line(fitz.Point(50, 150), fitz.Point(250, 150))
    page.draw_line(fitz.Point(50, 50), fitz.Point(50, 150))
    page.draw_line(fitz.Point(150, 50), fitz.Point(150, 150))
    page.draw_line(fitz.Point(250, 50), fitz.Point(250, 150))
    page.insert_text((60, 70), "A")
    page.insert_text((160, 70), "B")
    page.insert_text((60, 120), "C")
    page.insert_text((160, 120), "D")
    working = session_dir / "working.pdf"
    doc.save(working)
    doc.close()
    session = DocumentSession(working_path=working, source_path=None)

    result = session.apply(PdfToXlsxOperation())
    assert result.working_path.suffix == ".xlsx"
    workbook = openpyxl.load_workbook(result.working_path)
    all_values = {cell.value for sheet in workbook.worksheets for row in sheet.iter_rows() for cell in row}
    assert {"A", "B", "C", "D"} & {v for v in all_values if v}


def test_pdf_to_xlsx_with_no_tables_still_produces_valid_workbook(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("no tables here", None)])
    result = session.apply(PdfToXlsxOperation())
    workbook = openpyxl.load_workbook(result.working_path)
    assert workbook.worksheets


# --- PdfToHtmlOperation -----------------------------------------------------


def test_pdf_to_html_contains_escaped_text(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("<script>alert(1)</script>", None)])
    result = session.apply(PdfToHtmlOperation())
    html_text = result.working_path.read_text(encoding="utf-8")
    assert "<script>" not in html_text
    assert "&lt;script&gt;" in html_text
    assert 'class="page"' in html_text


def test_pdf_to_html_one_page_div_per_source_page(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("p1", None), ("p2", None)])
    result = session.apply(PdfToHtmlOperation())
    html_text = result.working_path.read_text(encoding="utf-8")
    assert html_text.count('class="page"') == 2


# --- PdfToJpgOperation -------------------------------------------------------


def test_pdf_to_jpg_renders_requested_page(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("p1", (1, 0, 0)), ("p2", (0, 1, 0))])
    result = session.apply(PdfToJpgOperation(page=2, dpi=100))
    assert result.working_path.suffix == ".jpg"
    assert result.working_path.stat().st_size > 0

    with fitz.open(result.working_path) as img:
        assert img.page_count == 1


def test_pdf_to_jpg_rejects_out_of_range_page(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("p1", None)])
    with pytest.raises(OperationError):
        session.apply(PdfToJpgOperation(page=5))


def test_pdf_to_jpg_rejects_too_low_dpi() -> None:
    with pytest.raises(OperationError):
        PdfToJpgOperation(page=1, dpi=10)


def test_pdf_to_jpg_undo_restores_pdf(tmp_path: Path) -> None:
    session = _session_with_pdf(tmp_path, [("p1", None)])
    result = session.apply(PdfToJpgOperation(page=1))
    restored = result.undo()
    assert restored.working_path.suffix == ".pdf"
