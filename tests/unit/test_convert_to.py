"""Unit tests for core/ops/convert_to.py (Word/PowerPoint/Excel/
HTML/JPG -> PDF). The four Office/HTML operations are dual-engine
(LibreOffice primary, pure-Python fallback) - every test forces the
fallback via monkeypatching `libreoffice_binary` so the suite stays
deterministic across machines that may not have `soffice` installed,
plus one skipif-guarded test per operation that exercises the real
LibreOffice subprocess when available. JpgToPdfOperation is fitz-only,
no engine split needed."""

from __future__ import annotations

from pathlib import Path

import docx
import fitz
import openpyxl
import pdfplumber
import pikepdf
import pytest
from pptx import Presentation
from pptx.util import Emu

from core.errors import OperationError
from core.model.document import DocumentSession
from core.ops.convert_common import libreoffice_binary
from core.ops.convert_to import (
    DocxToPdfOperation,
    HtmlToPdfOperation,
    JpgToPdfOperation,
    PptxToPdfOperation,
    XlsxToPdfOperation,
)

_HAS_LIBREOFFICE = libreoffice_binary() is not None


@pytest.fixture
def force_fallback(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr("core.ops.convert_to.libreoffice_binary", lambda: None)


def _fresh_session(tmp_path: Path, name: str = "session") -> DocumentSession:
    session_dir = tmp_path / name
    session_dir.mkdir(exist_ok=True)
    return DocumentSession(working_path=session_dir / "placeholder.pdf", source_path=None)


def _make_docx(path: Path) -> Path:
    document = docx.Document()
    document.add_heading("Title", level=1)
    document.add_paragraph("Hello from a test docx.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "C"
    table.cell(1, 1).text = "D"
    document.save(str(path))
    return path


def _make_pptx(path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height // 4)
    textbox.text_frame.text = "Slide text"
    prs.save(str(path))
    return path


def _make_xlsx(path: Path) -> Path:
    workbook = openpyxl.Workbook()
    workbook.active.append(["col1", "col2"])
    workbook.active.append([1, 2])
    workbook.save(str(path))
    return path


def _make_html(path: Path, body: str = "<h1>Hi</h1><p>paragraph</p>") -> Path:
    path.write_text(f"<html><body>{body}</body></html>", encoding="utf-8")
    return path


def _make_docx_with_a_table_between_paragraphs(path: Path) -> Path:
    """A body whose children interleave - the shape that catches a
    converter reading all paragraphs before all tables."""
    document = docx.Document()
    document.add_paragraph("PARA-BEFORE-TABLE")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "TABLE-CELL"
    document.add_paragraph("PARA-AFTER-TABLE")
    document.save(str(path))
    return path


def _page_count(pdf_path: Path) -> int:
    with pikepdf.Pdf.open(pdf_path) as pdf:
        return len(pdf.pages)


def _pdf_text(pdf_path: Path) -> str:
    with fitz.open(str(pdf_path)) as pdf:
        return "\n".join(page.get_text() for page in pdf)


# --- DocxToPdfOperation ----------------------------------------------------


def test_docx_to_pdf_fallback_produces_valid_pdf(tmp_path: Path, force_fallback: None) -> None:
    source = _make_docx(tmp_path / "in.docx")
    session = _fresh_session(tmp_path)
    result = session.apply(DocxToPdfOperation(source_path=source))
    assert result.working_path.suffix == ".pdf"
    assert _page_count(result.working_path) >= 1
    assert "pure-Python fallback" in result.operation_log[-1].describe()


def test_docx_to_pdf_missing_source_raises() -> None:
    doc = DocumentSession(working_path=None, source_path=None)
    with pytest.raises(OperationError):
        doc.apply(DocxToPdfOperation(source_path=Path("/nonexistent/file.docx")))


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_docx_to_pdf_via_libreoffice(tmp_path: Path) -> None:
    source = _make_docx(tmp_path / "in.docx")
    session = _fresh_session(tmp_path)
    result = session.apply(DocxToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "LibreOffice" in result.operation_log[-1].describe()


# The two tests above check that a PDF came out and which engine made
# it, but never that the Word document's own content survived - a
# converter that produced a blank page would pass both. These assert
# the text actually lands in the output, for each engine separately.


def test_docx_to_pdf_fallback_preserves_the_documents_text(
    tmp_path: Path, force_fallback: None
) -> None:
    source = _make_docx(tmp_path / "in.docx")
    session = _fresh_session(tmp_path)
    result = session.apply(DocxToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    assert "Title" in text
    assert "Hello from a test docx." in text
    for cell in ("A", "B", "C", "D"):
        assert cell in text


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_docx_to_pdf_via_libreoffice_preserves_the_documents_text(tmp_path: Path) -> None:
    source = _make_docx(tmp_path / "in.docx")
    session = _fresh_session(tmp_path)
    result = session.apply(DocxToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    assert "Title" in text
    assert "Hello from a test docx." in text
    for cell in ("A", "B", "C", "D"):
        assert cell in text


def test_docx_to_pdf_fallback_keeps_paragraphs_and_tables_in_document_order(
    tmp_path: Path, force_fallback: None
) -> None:
    """Regression: the fallback used to walk `document.paragraphs` and
    then `document.tables`, two collections that each flatten the body
    separately - so a table sitting between two paragraphs came out
    after both of them, silently reordering the document's content."""
    source = _make_docx_with_a_table_between_paragraphs(tmp_path / "order.docx")
    session = _fresh_session(tmp_path)
    result = session.apply(DocxToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    assert (
        text.index("PARA-BEFORE-TABLE") < text.index("TABLE-CELL") < text.index("PARA-AFTER-TABLE")
    )


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_both_docx_engines_agree_on_content_order(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    """The fallback is a lower-fidelity reconstruction, but it must not
    disagree with the primary engine about what order the content is
    in - that's the one thing a text reconstruction has to get right.
    Runs both engines on the same document rather than using the
    module's `force_fallback` fixture, which is all-or-nothing."""
    source = _make_docx_with_a_table_between_paragraphs(tmp_path / "order.docx")
    markers = ("PARA-BEFORE-TABLE", "TABLE-CELL", "PARA-AFTER-TABLE")

    libreoffice_result = _fresh_session(tmp_path, "lo").apply(
        DocxToPdfOperation(source_path=source)
    )
    assert "LibreOffice" in libreoffice_result.operation_log[-1].describe()
    libreoffice_text = _pdf_text(libreoffice_result.working_path)

    monkeypatch.setattr("core.ops.convert_to.libreoffice_binary", lambda: None)
    fallback_result = _fresh_session(tmp_path, "fb").apply(DocxToPdfOperation(source_path=source))
    assert "pure-Python fallback" in fallback_result.operation_log[-1].describe()
    fallback_text = _pdf_text(fallback_result.working_path)

    assert [m for m in markers if m in libreoffice_text] == list(markers)
    assert sorted(markers, key=libreoffice_text.index) == sorted(markers, key=fallback_text.index)


# --- PptxToPdfOperation ------------------------------------------------------


def test_pptx_to_pdf_fallback_produces_valid_pdf(tmp_path: Path, force_fallback: None) -> None:
    source = _make_pptx(tmp_path / "in.pptx")
    session = _fresh_session(tmp_path)
    result = session.apply(PptxToPdfOperation(source_path=source))
    assert _page_count(result.working_path) == 1
    assert "pure-Python fallback" in result.operation_log[-1].describe()


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_pptx_to_pdf_via_libreoffice(tmp_path: Path) -> None:
    source = _make_pptx(tmp_path / "in.pptx")
    session = _fresh_session(tmp_path)
    result = session.apply(PptxToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "LibreOffice" in result.operation_log[-1].describe()


# --- XlsxToPdfOperation -------------------------------------------------------


def test_xlsx_to_pdf_fallback_produces_valid_pdf(tmp_path: Path, force_fallback: None) -> None:
    source = _make_xlsx(tmp_path / "in.xlsx")
    session = _fresh_session(tmp_path)
    result = session.apply(XlsxToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "pure-Python fallback" in result.operation_log[-1].describe()


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_xlsx_to_pdf_via_libreoffice(tmp_path: Path) -> None:
    source = _make_xlsx(tmp_path / "in.xlsx")
    session = _fresh_session(tmp_path)
    result = session.apply(XlsxToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "LibreOffice" in result.operation_log[-1].describe()


# --- HtmlToPdfOperation -------------------------------------------------------


def test_html_to_pdf_fallback_produces_valid_pdf(tmp_path: Path, force_fallback: None) -> None:
    source = _make_html(tmp_path / "in.html")
    session = _fresh_session(tmp_path)
    result = session.apply(HtmlToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "xhtml2pdf" in result.operation_log[-1].describe()


def test_html_to_pdf_fallback_rejects_remote_resources(tmp_path: Path, force_fallback: None) -> None:
    source = _make_html(tmp_path / "in.html", body='<img src="http://evil.example/x.png">')
    session = _fresh_session(tmp_path)
    with pytest.raises(OperationError, match="[Rr]efus"):
        session.apply(HtmlToPdfOperation(source_path=source))


def test_html_to_pdf_fallback_allows_data_uri_images(tmp_path: Path, force_fallback: None) -> None:
    # 1x1 transparent PNG
    data_uri = (
        "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAAB"
        "CAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
    )
    source = _make_html(tmp_path / "in.html", body=f'<img src="{data_uri}">')
    session = _fresh_session(tmp_path)
    result = session.apply(HtmlToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed on this machine")
def test_html_to_pdf_via_libreoffice(tmp_path: Path) -> None:
    source = _make_html(tmp_path / "in.html")
    session = _fresh_session(tmp_path)
    result = session.apply(HtmlToPdfOperation(source_path=source))
    assert _page_count(result.working_path) >= 1
    assert "LibreOffice" in result.operation_log[-1].describe()


# --- JpgToPdfOperation ------------------------------------------------------


def _make_jpg(path: Path, color: tuple[float, float, float]) -> Path:
    doc = fitz.open()
    page = doc.new_page(width=120, height=80)
    page.draw_rect(fitz.Rect(0, 0, 60, 40), color=color, fill=color)
    page.get_pixmap().save(str(path))
    doc.close()
    return path


def test_jpg_to_pdf_combines_images_into_one_pdf(tmp_path: Path) -> None:
    img1 = _make_jpg(tmp_path / "img1.jpg", (1, 0, 0))
    img2 = _make_jpg(tmp_path / "img2.jpg", (0, 1, 0))
    session = _fresh_session(tmp_path)
    result = session.apply(JpgToPdfOperation(sources=[img1, img2]))
    assert _page_count(result.working_path) == 2


def test_jpg_to_pdf_requires_at_least_one_source(tmp_path: Path) -> None:
    session = _fresh_session(tmp_path)
    with pytest.raises(OperationError):
        session.apply(JpgToPdfOperation(sources=[]))


def test_jpg_to_pdf_missing_source_raises(tmp_path: Path) -> None:
    session = _fresh_session(tmp_path)
    with pytest.raises(OperationError):
        session.apply(JpgToPdfOperation(sources=[Path("/nonexistent/img.jpg")]))


def test_jpg_to_pdf_undo_restores_no_document_state(tmp_path: Path) -> None:
    # _fresh_session's placeholder.pdf never actually exists (there is
    # no real prior document, only a directory-resolution stand-in -
    # see cli/main.py's identical pattern), so the pre-apply snapshot
    # is None and undo should restore to "nothing open," not conjure a
    # document that was never really there.
    img1 = _make_jpg(tmp_path / "img1.jpg", (1, 0, 0))
    session = _fresh_session(tmp_path)
    result = session.apply(JpgToPdfOperation(sources=[img1]))
    restored = result.undo()
    assert restored.working_path is None


# --- content, not just page counts -------------------------------------------
#
# Every test above this line checks page count and engine name only, so
# a conversion that emitted blank pages passed them all. These assert
# the source file's content actually arrives in the PDF, and that it
# arrives somewhere a reader can actually see - text sitting off the
# edge of the page still shows up in extraction.


def _words_off_page(pdf_path: Path) -> list[str]:
    """Words whose bounding box escapes the page. Text placed past the
    right edge is invisible when the PDF is viewed or printed, but is
    still found by text extraction - so a plain "is the text present"
    assertion cannot catch it.

    pdfplumber rather than fitz deliberately: fitz's `get_text("words")`
    did not report the overflow for a single very long cell value that
    pdfplumber measured at x1=2980 on a 612pt page, so a fitz-based
    version of this helper silently passed against the unfixed code.
    """
    escaped: list[str] = []
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            for word in page.extract_words():
                if (
                    word["x1"] > page.width + 0.5
                    or word["x0"] < -0.5
                    or word["bottom"] > page.height + 0.5
                ):
                    escaped.append(str(word["text"]))
    return escaped


def _make_pptx_with_a_table(path: Path) -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Emu(457200), Emu(457200), Emu(4572000), Emu(1828800))
    for (row, column), text in {
        (0, 0): "TABLE-CELL-A",
        (0, 1): "TABLE-CELL-B",
        (1, 0): "TABLE-CELL-C",
        (1, 1): "TABLE-CELL-D",
    }.items():
        shape.table.cell(row, column).text = text
    prs.save(str(path))
    return path


def test_pptx_to_pdf_fallback_preserves_slide_text(tmp_path: Path, force_fallback: None) -> None:
    source = _make_pptx(tmp_path / "in.pptx")
    session = _fresh_session(tmp_path)
    result = session.apply(PptxToPdfOperation(source_path=source))
    assert "Slide text" in _pdf_text(result.working_path)


def test_pptx_to_pdf_fallback_renders_table_shapes(tmp_path: Path, force_fallback: None) -> None:
    """Regression: table shapes were skipped entirely by the fallback,
    so a slide whose content was a table converted to a blank page -
    while LibreOffice rendered it. A table is text, which is exactly
    what this reconstruction exists to preserve."""
    source = _make_pptx_with_a_table(tmp_path / "table.pptx")
    session = _fresh_session(tmp_path)
    result = session.apply(PptxToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    for cell in ("TABLE-CELL-A", "TABLE-CELL-B", "TABLE-CELL-C", "TABLE-CELL-D"):
        assert cell in text


def test_pptx_to_pdf_fallback_wraps_text_inside_the_slide(
    tmp_path: Path, force_fallback: None
) -> None:
    """Regression: `Canvas.drawString` does no wrapping, so a text box
    longer than its width was drawn as one unbroken line running off
    the side of the slide - extractable, but invisible to a reader."""
    path = tmp_path / "long.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Emu(457200), Emu(457200), Emu(1828800), Emu(914400))
    box.text_frame.text = "LONGSTART " + ("wrap " * 120) + "LONGEND"
    prs.save(str(path))

    session = _fresh_session(tmp_path)
    result = session.apply(PptxToPdfOperation(source_path=path))
    text = _pdf_text(result.working_path)
    assert "LONGSTART" in text and "LONGEND" in text
    assert _words_off_page(result.working_path) == []


def test_xlsx_to_pdf_fallback_preserves_cell_values(tmp_path: Path, force_fallback: None) -> None:
    source = _make_xlsx(tmp_path / "in.xlsx")
    session = _fresh_session(tmp_path)
    result = session.apply(XlsxToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    for value in ("col1", "col2", "1", "2"):
        assert value in text


def test_xlsx_to_pdf_fallback_fits_a_wide_sheet_on_the_page(
    tmp_path: Path, force_fallback: None
) -> None:
    """Regression: a bare `Table(data)` sizes columns to their widest
    cell with no upper bound, so a 25-column sheet ran off both edges
    of the page (measured: 12 of 51 words escaped)."""
    path = tmp_path / "wide.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append([f"COL{i:02d}" for i in range(25)])
    sheet.append([f"val{i:02d}" for i in range(25)])
    workbook.save(str(path))

    session = _fresh_session(tmp_path)
    result = session.apply(XlsxToPdfOperation(source_path=path))
    assert "COL00" in _pdf_text(result.working_path)
    assert _words_off_page(result.working_path) == []


def test_xlsx_to_pdf_fallback_fits_a_long_cell_value_on_the_page(
    tmp_path: Path, force_fallback: None
) -> None:
    """The same overflow from a single unbroken value rather than from
    column count - it has no space to wrap at, hence splitLongWords."""
    path = tmp_path / "longcell.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.append(["A" * 400, "B" * 400])
    workbook.save(str(path))

    session = _fresh_session(tmp_path)
    result = session.apply(XlsxToPdfOperation(source_path=path))
    assert _words_off_page(result.working_path) == []


def test_html_to_pdf_fallback_preserves_content(tmp_path: Path, force_fallback: None) -> None:
    source = _make_html(
        tmp_path / "in.html", body="<h1>HTML-HEADING</h1><p>HTML-PARAGRAPH</p>"
    )
    session = _fresh_session(tmp_path)
    result = session.apply(HtmlToPdfOperation(source_path=source))
    text = _pdf_text(result.working_path)
    assert "HTML-HEADING" in text
    assert "HTML-PARAGRAPH" in text


def test_jpg_to_pdf_really_embeds_each_image_in_order(tmp_path: Path) -> None:
    """The existing coverage asserted a two-page PDF came out, which a
    pair of blank pages would satisfy. This checks the images are
    genuinely embedded, and that page order follows source order."""
    red = _make_jpg(tmp_path / "red.jpg", (1, 0, 0))
    green = _make_jpg(tmp_path / "green.jpg", (0, 1, 0))
    session = _fresh_session(tmp_path)
    result = session.apply(JpgToPdfOperation(sources=[red, green]))

    with fitz.open(str(result.working_path)) as pdf:
        assert len(pdf) == 2
        assert [len(pdf[i].get_images()) for i in range(2)] == [1, 1]
        samples = []
        for page in pdf:
            pixmap = page.get_pixmap()
            # _make_jpg fills only the top-left quadrant, so the page
            # centre sits exactly on the colour boundary and samples a
            # JPEG-blurred blend. A quarter in is safely inside it.
            samples.append(pixmap.pixel(pixmap.width // 4, pixmap.height // 4))
    red_page, green_page = samples
    assert red_page[0] > 200 and red_page[1] < 60
    assert green_page[1] > 200 and green_page[0] < 60
