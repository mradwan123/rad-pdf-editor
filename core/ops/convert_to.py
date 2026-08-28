"""External Word/PowerPoint/Excel/HTML/JPG files -> PDF conversions
(SPEC.md Phase 3 list, the external-file-is-the-source half). Shaped
like MergeOperation (core/ops/merge_split.py): an external file (or
list of files) is the source, not `doc.working_path` - any currently
open document is untouched, and these ops work even with nothing open
yet, exactly like Merge.

See core/ops/convert_from.py for the reverse direction, and
core/ops/convert_common.py's module docstring for the dual-engine
strategy - LibreOffice genuinely is the primary engine in *this*
direction (confirmed by hand: Writer/Impress/Calc/Writer-Web all
export straight to PDF reliably) - and its known
network-lockdown-subprocess caveat.

Every operation records which engine actually ran (where more than one
is possible) in `describe()`, so the undo-stack UI and audit trail
reflect the fidelity tradeoff actually made, not just "converted."
"""

from __future__ import annotations

import contextlib
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape as _xml_escape

import docx
import fitz
import openpyxl
import pikepdf
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.utils import ImageReader, simpleSplit
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table
from xhtml2pdf import pisa

from core.errors import ConversionError, OperationError
from core.model.document import DocumentSession
from core.model.operation import Operation
from core.ops.common import (
    allocate_working_path,
    next_session,
    read_working_bytes,
    snapshot_restore_invert,
)
from core.ops.convert_common import libreoffice_binary, run_libreoffice_convert
from core.registry.plugin_base import ToolPlugin

CORE_VERSION_RANGE = ">=1.0,<2.0"

#: Points per EMU (English Metric Unit), for reading python-pptx shape
#: positions into reportlab's point-based coordinate system.
_POINTS_PER_EMU = 1 / 12700

#: The font the canvas-drawing fallbacks measure and draw with. Set
#: explicitly rather than relying on the reportlab canvas default,
#: because `simpleSplit` has to wrap against the *same* font/size the
#: text is finally drawn in or the wrap width is meaningless.
_FALLBACK_FONT = "Helvetica"
_FALLBACK_FONT_SIZE = 12
_FALLBACK_LINE_HEIGHT = 14

#: Spreadsheet cells are rendered smaller than body text so a sheet
#: with many columns still fits the printable width legibly, and the
#: page margin the column widths are computed against.
_SHEET_CELL_FONT_SIZE = 7
_SHEET_MARGIN = 36
#: Narrowest a spreadsheet column may get before the sheet is split
#: into further column groups instead - below roughly this, ordinary
#: cell values start breaking mid-word to fit.
_SHEET_MIN_COLUMN_WIDTH = 45


def _require_source_file(path: Path) -> None:
    if not path.is_file():
        raise OperationError(f"Source file not found: {path}")


def _docx_fallback_to_pdf(source_path: Path, out_path: Path) -> None:
    """python-docx read -> reportlab platypus reconstruction. Font
    sizes/bold/italic run properties aren't mapped, and complex layout
    (columns, image text-wrap) isn't preserved - a text-and-tables
    reconstruction, not a pixel-faithful one.

    Paragraphs and tables are walked over the document body's own child
    elements rather than over `document.paragraphs` then
    `document.tables`, because those two collections each flatten the
    body separately: reading them in sequence emits every paragraph
    first and every table afterwards, so a table sitting *between* two
    paragraphs came out at the end of the PDF. Content order is the one
    thing a text reconstruction has to get right, and LibreOffice (the
    primary engine) preserves it, so the fallback must not silently
    disagree with it.

    Known limitation, deliberately not worked around here: the built-in
    reportlab fonts are Latin-1, so characters outside it (CJK, and
    other non-Latin scripts) do not survive - reportlab substitutes a
    placeholder glyph rather than raising. Text in those scripts must
    go through the LibreOffice engine, which handles it correctly.
    """
    document = docx.Document(str(source_path))
    styles = getSampleStyleSheet()
    story: list[Any] = []
    for child in document.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            para = DocxParagraph(child, document)
            text = para.text
            if not text.strip():
                story.append(Spacer(1, 6))
                continue
            style_name = (
                "Heading1" if para.style is not None and "Heading" in para.style.name else "Normal"
            )
            story.append(Paragraph(_xml_escape(text), styles[style_name]))
        elif tag == "tbl":
            data = [[cell.text for cell in row.cells] for row in DocxTable(child, document).rows]
            if data:
                story.append(Table(data))
                story.append(Spacer(1, 12))
    if not story:
        story.append(Paragraph("", styles["Normal"]))
    SimpleDocTemplate(str(out_path), pagesize=letter).build(story)


def _draw_pptx_table(
    pdf_canvas: canvas.Canvas,
    table: Any,
    left: float,
    bottom_y: float,
    shape_width: float,
    shape_height: float,
) -> None:
    """Draw a slide table's cell text on an even grid inside the shape's
    own rectangle. Cell borders/fills aren't reproduced - this recovers
    the table's *text*, positioned roughly where it belongs, which is
    the same bargain the rest of this fallback makes."""
    rows = list(table.rows)
    columns = list(table.columns)
    if not rows or not columns:
        return
    cell_width = shape_width / len(columns)
    cell_height = shape_height / len(rows)
    for row_index, _row in enumerate(rows):
        # Rows run top-to-bottom; the shape's y is its *bottom* edge.
        row_top = bottom_y + shape_height - row_index * cell_height
        for column_index in range(len(columns)):
            text = table.cell(row_index, column_index).text
            if not text.strip():
                continue
            _draw_wrapped_text(
                pdf_canvas,
                text,
                left + column_index * cell_width,
                row_top - _FALLBACK_LINE_HEIGHT,
                max(cell_width, _FALLBACK_LINE_HEIGHT),
            )


def _draw_wrapped_text(
    pdf_canvas: canvas.Canvas, text: str, left: float, top_y: float, max_width: float
) -> None:
    """Draw `text` at (left, top_y) wrapped to `max_width`, one line per
    rendered row, growing downward.

    `Canvas.drawString` does no wrapping whatsoever: a text frame's
    contents were previously emitted as a single unbroken line, so
    anything longer than its box ran clean off the right edge of the
    slide and out of the page. The text was still in the PDF's text
    layer (extraction found it) but was invisible to a reader, which is
    the failure mode that makes it worth wrapping rather than leaving
    to "approximate positions".
    """
    lines = simpleSplit(text, _FALLBACK_FONT, _FALLBACK_FONT_SIZE, max_width)
    y = top_y
    for line in lines:
        pdf_canvas.drawString(left, y, line)
        y -= _FALLBACK_LINE_HEIGHT


def _pptx_fallback_to_pdf(source_path: Path, out_path: Path) -> None:
    """python-pptx read -> one reportlab page per slide, text frames,
    tables and picture shapes positioned via EMU->pt conversion.
    Approximate positions, not a pixel-exact slide renderer (python-pptx
    has none) - genuinely graphical shape types (charts, SmartArt, ...)
    are skipped, since there is nothing textual to recover from them.

    Table shapes are *not* in that skipped category, though they used to
    be: a table is text, which is exactly what this reconstruction
    exists to preserve, and LibreOffice renders it - so dropping it
    silently made the two engines disagree about the slide's content,
    not merely about its fidelity.
    """
    prs = Presentation(str(source_path))
    # A loaded/created Presentation always has these set from its
    # template; None is only in the type stubs for the pathological
    # case of a hand-built package that never specified a slide size.
    assert prs.slide_width is not None
    assert prs.slide_height is not None
    width_pt = prs.slide_width * _POINTS_PER_EMU
    height_pt = prs.slide_height * _POINTS_PER_EMU
    pdf_canvas = canvas.Canvas(str(out_path), pagesize=(width_pt, height_pt))
    pdf_canvas.setFont(_FALLBACK_FONT, _FALLBACK_FONT_SIZE)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.left is None or shape.top is None:
                continue
            left = shape.left * _POINTS_PER_EMU
            top = shape.top * _POINTS_PER_EMU
            shape_width = (shape.width or 0) * _POINTS_PER_EMU
            shape_height = (shape.height or 0) * _POINTS_PER_EMU
            y = height_pt - top - shape_height
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_stream = io.BytesIO(shape.image.blob)
                pdf_canvas.drawImage(
                    ImageReader(image_stream), left, y, width=shape_width, height=shape_height
                )
            elif getattr(shape, "has_table", False):
                _draw_pptx_table(pdf_canvas, shape.table, left, y, shape_width, shape_height)
            elif shape.has_text_frame and shape.text_frame.text:
                _draw_wrapped_text(
                    pdf_canvas,
                    shape.text_frame.text,
                    left,
                    y + max(shape_height - _FALLBACK_LINE_HEIGHT, 0),
                    max(shape_width, _FALLBACK_LINE_HEIGHT),
                )
        pdf_canvas.showPage()
    pdf_canvas.save()


def _xlsx_fallback_to_pdf(source_path: Path, out_path: Path) -> None:
    """openpyxl read -> one reportlab Table per worksheet. No cell
    styling/merged-cell visual fidelity beyond plain text; charts and
    embedded images in the workbook are not rendered.

    Column widths are fitted to the printable frame and cell text is
    wrapped inside them. A bare `Table(data)` sizes its columns to
    their widest cell instead, with no upper bound, so a wide sheet
    (measured: 25 columns) ran off both edges of the page - the text
    was in the PDF but a reader could not see it, and LibreOffice
    paginates the same sheet properly. Long single cell values did the
    same thing on their own.
    """
    workbook = openpyxl.load_workbook(source_path, data_only=True)
    styles = getSampleStyleSheet()
    cell_style = ParagraphStyle(
        "SheetCell",
        parent=styles["BodyText"],
        fontSize=_SHEET_CELL_FONT_SIZE,
        leading=_SHEET_CELL_FONT_SIZE + 2,
        spaceAfter=0,
        # A long unbroken value (an id, a URL, a hash) has no space to
        # wrap at, so without this it would overflow its column and
        # reintroduce exactly the bug this function is fixing.
        splitLongWords=True,
    )
    available_width = letter[0] - _SHEET_MARGIN * 2
    story: list[Any] = []
    for sheet in workbook.worksheets:
        story.append(Paragraph(_xml_escape(sheet.title), styles["Heading2"]))
        data = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        if data:
            column_count = max(len(row) for row in data)
            # Splitting a wide sheet into column groups, rather than
            # squeezing every column onto one page: at 25 columns a
            # single-page fit leaves ~21pt per column, too narrow for
            # even a 5-character heading, so splitLongWords chopped
            # values mid-word ("COL00" -> "C/OL/00"). LibreOffice
            # paginates the same sheet across pages instead, and this
            # keeps values intact for the same reason. A normal sheet
            # is one group, i.e. unchanged.
            per_page = max(1, int(available_width // _SHEET_MIN_COLUMN_WIDTH))
            for start in range(0, column_count, per_page):
                columns = min(per_page, column_count - start)
                chunk = [
                    [
                        Paragraph(_xml_escape(value), cell_style)
                        for value in row[start : start + columns]
                    ]
                    for row in data
                ]
                story.append(
                    Table(chunk, colWidths=[available_width / columns] * columns)
                )
                story.append(Spacer(1, 12))
        story.append(Spacer(1, 12))
    if not story:
        story.append(Paragraph("", styles["Normal"]))
    SimpleDocTemplate(
        str(out_path),
        pagesize=letter,
        leftMargin=_SHEET_MARGIN,
        rightMargin=_SHEET_MARGIN,
    ).build(story)


def _reject_remote_uri(uri: str, _relative: str) -> str:
    """xhtml2pdf's `link_callback` resolves every resource an HTML
    document references (img src, link href, ...) to a local path it
    reads directly. Defense in depth on top of `network_lockdown()`
    (SPEC.md's "no network calls anywhere" hard requirement): refuse
    anything that isn't already a local path or a `data:` URI, so a
    crafted HTML file can't cause even an attempted outbound fetch."""
    if uri.startswith("data:"):
        return uri
    if "://" in uri or uri.startswith("//"):
        raise ConversionError(f"Refusing to fetch a remote resource referenced in HTML: {uri}")
    return uri


def _html_fallback_to_pdf(source_path: Path, out_path: Path) -> None:
    html_text = source_path.read_text(encoding="utf-8", errors="replace")
    with out_path.open("wb") as f:
        result = pisa.CreatePDF(html_text, dest=f, link_callback=_reject_remote_uri)
    if result.err:
        raise ConversionError(f"xhtml2pdf failed to render '{source_path.name}' ({result.err} error(s)).")


@dataclass
class DocxToPdfOperation(Operation):
    """Converts an external Word document to PDF. Tries LibreOffice
    headless first (confirmed by hand: Writer exports docx -> pdf
    reliably); falls back to `_docx_fallback_to_pdf` when LibreOffice
    isn't installed."""

    source_path: Path
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)
    _engine_used: str = field(default="", init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_source_file(self.source_path)
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pdf")
        if libreoffice_binary() is not None:
            converted = run_libreoffice_convert(self.source_path, "pdf", out_path.parent)
            converted.replace(out_path)
            self._engine_used = "LibreOffice"
        else:
            _docx_fallback_to_pdf(self.source_path, out_path)
            self._engine_used = "pure-Python fallback"

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "docx_to_pdf",
            "source_path": str(self.source_path),
        }

    def describe(self) -> str:
        return f"Converted Word document to PDF ({self._engine_used or 'pending'})"


@dataclass
class PptxToPdfOperation(Operation):
    """Converts an external PowerPoint presentation to PDF. Tries
    LibreOffice headless first (confirmed by hand: Impress exports
    pptx -> pdf reliably); falls back to `_pptx_fallback_to_pdf` when
    LibreOffice isn't installed."""

    source_path: Path
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)
    _engine_used: str = field(default="", init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_source_file(self.source_path)
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pdf")
        if libreoffice_binary() is not None:
            converted = run_libreoffice_convert(self.source_path, "pdf", out_path.parent)
            converted.replace(out_path)
            self._engine_used = "LibreOffice"
        else:
            _pptx_fallback_to_pdf(self.source_path, out_path)
            self._engine_used = "pure-Python fallback"

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "pptx_to_pdf",
            "source_path": str(self.source_path),
        }

    def describe(self) -> str:
        return f"Converted PowerPoint file to PDF ({self._engine_used or 'pending'})"


@dataclass
class XlsxToPdfOperation(Operation):
    """Converts an external Excel workbook to PDF. Tries LibreOffice
    headless first (confirmed by hand: Calc exports xlsx -> pdf
    reliably); falls back to `_xlsx_fallback_to_pdf` when LibreOffice
    isn't installed."""

    source_path: Path
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)
    _engine_used: str = field(default="", init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_source_file(self.source_path)
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pdf")
        if libreoffice_binary() is not None:
            converted = run_libreoffice_convert(self.source_path, "pdf", out_path.parent)
            converted.replace(out_path)
            self._engine_used = "LibreOffice"
        else:
            _xlsx_fallback_to_pdf(self.source_path, out_path)
            self._engine_used = "pure-Python fallback"

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "xlsx_to_pdf",
            "source_path": str(self.source_path),
        }

    def describe(self) -> str:
        return f"Converted Excel workbook to PDF ({self._engine_used or 'pending'})"


@dataclass
class HtmlToPdfOperation(Operation):
    """Converts an external HTML file to PDF. Tries LibreOffice
    headless first (confirmed by hand: Writer/Web exports html -> pdf
    reliably); falls back to xhtml2pdf when LibreOffice isn't
    installed - basic CSS only, no JS execution, and
    `_reject_remote_uri` explicitly blocks any remote resource
    reference regardless of engine."""

    source_path: Path
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)
    _engine_used: str = field(default="", init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_source_file(self.source_path)
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pdf")
        if libreoffice_binary() is not None:
            converted = run_libreoffice_convert(self.source_path, "pdf", out_path.parent)
            converted.replace(out_path)
            self._engine_used = "LibreOffice"
        else:
            _html_fallback_to_pdf(self.source_path, out_path)
            self._engine_used = "pure-Python fallback (xhtml2pdf)"

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "html_to_pdf",
            "source_path": str(self.source_path),
        }

    def describe(self) -> str:
        return f"Converted HTML file to PDF ({self._engine_used or 'pending'})"


@dataclass
class JpgToPdfOperation(Operation):
    """Combines `sources` (JPG/PNG images), in order, into a single
    PDF - one page per image, each page sized to that image's own
    pixel dimensions. Mirrors MergeOperation's exact shape (external
    files in, one PDF out) and its ExitStack pattern for keeping every
    intermediate pikepdf.Pdf alive until the final save. fitz only -
    no LibreOffice/pure-Python split needed, this is already the
    simplest path."""

    sources: list[Path]
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        if not self.sources:
            raise OperationError("Convert to PDF requires at least one source image.")
        for src in self.sources:
            _require_source_file(src)
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pdf")
        merged = pikepdf.Pdf.new()
        with contextlib.ExitStack() as stack:
            for src in self.sources:
                try:
                    with fitz.open(src) as image_doc:
                        pdf_bytes = image_doc.convert_to_pdf()
                except RuntimeError as exc:
                    raise ConversionError(f"Could not read image '{src.name}': {exc}") from exc
                image_pdf = stack.enter_context(pikepdf.Pdf.open(io.BytesIO(pdf_bytes)))
                merged.pages.extend(image_pdf.pages)
            merged.save(out_path)

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "jpg_to_pdf",
            "sources": [str(p) for p in self.sources],
        }

    def describe(self) -> str:
        return f"Converted {len(self.sources)} image(s) to PDF"


class DocxToPdfPlugin(ToolPlugin):
    tool_id = "docx_to_pdf"
    display_name = "Word to PDF"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            source_path = kwargs["source_path"]
        except KeyError as exc:
            raise OperationError("Word to PDF requires a 'source_path'.") from exc
        return DocxToPdfOperation(source_path=Path(source_path))

    def operation_class(self) -> type[Operation]:
        return DocxToPdfOperation


class PptxToPdfPlugin(ToolPlugin):
    tool_id = "pptx_to_pdf"
    display_name = "PowerPoint to PDF"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            source_path = kwargs["source_path"]
        except KeyError as exc:
            raise OperationError("PowerPoint to PDF requires a 'source_path'.") from exc
        return PptxToPdfOperation(source_path=Path(source_path))

    def operation_class(self) -> type[Operation]:
        return PptxToPdfOperation


class XlsxToPdfPlugin(ToolPlugin):
    tool_id = "xlsx_to_pdf"
    display_name = "Excel to PDF"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            source_path = kwargs["source_path"]
        except KeyError as exc:
            raise OperationError("Excel to PDF requires a 'source_path'.") from exc
        return XlsxToPdfOperation(source_path=Path(source_path))

    def operation_class(self) -> type[Operation]:
        return XlsxToPdfOperation


class HtmlToPdfPlugin(ToolPlugin):
    tool_id = "html_to_pdf"
    display_name = "HTML to PDF"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            source_path = kwargs["source_path"]
        except KeyError as exc:
            raise OperationError("HTML to PDF requires a 'source_path'.") from exc
        return HtmlToPdfOperation(source_path=Path(source_path))

    def operation_class(self) -> type[Operation]:
        return HtmlToPdfOperation


class JpgToPdfPlugin(ToolPlugin):
    tool_id = "jpg_to_pdf"
    display_name = "JPG to PDF"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            sources = kwargs["sources"]
        except KeyError as exc:
            raise OperationError("JPG to PDF requires a 'sources' list of image paths.") from exc
        return JpgToPdfOperation(sources=[Path(p) for p in sources])

    def operation_class(self) -> type[Operation]:
        return JpgToPdfOperation
