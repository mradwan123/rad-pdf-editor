"""PDF -> Word/PowerPoint/Excel/HTML/JPG conversions (SPEC.md Phase 3
list, the PDF-is-the-source half). See core/ops/convert_to.py for the
external-file-is-the-source half (where LibreOffice genuinely is the
primary engine), and core/ops/convert_common.py's module docstring for
this project's network-lockdown-subprocess caveat.

Every operation in *this* module is pure-Python only - confirmed by
hand against the real `soffice` binary (not assumed) that LibreOffice
cannot do any "PDF -> Office format" conversion at all: a PDF always
imports as a *Draw* document, and Draw's export filter set only covers
odg/pdf/image formats, never docx/pptx ("Error: no export filter for
...sample.docx found, aborting." - reproduced for both docx and pptx
targets). Calc has no PDF-import filter whatsoever (a separate, even
more fundamental gap). So unlike convert_to.py, there is no working
LibreOffice path here to even attempt before falling back.
"""

from __future__ import annotations

import html
from dataclasses import dataclass, field
from typing import Any

import docx
import fitz
import openpyxl
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.util import Emu

from core.errors import OperationError
from core.model.document import DocumentSession
from core.model.operation import Operation
from core.ops.common import (
    allocate_working_path,
    next_session,
    read_working_bytes,
    snapshot_restore_invert,
)
from core.ops.convert_common import (
    extract_pdf_tables_by_page,
    extract_pdf_text_by_page,
    render_pdf_page_to_image,
)
from core.registry.plugin_base import ToolPlugin

CORE_VERSION_RANGE = ">=1.0,<2.0"

#: EMU (English Metric Units) per point - 1pt = 1/72in, 1in = 914400 EMU.
_EMU_PER_POINT = 12700


def _require_working_pdf(doc: DocumentSession) -> None:
    if doc.working_path is None:
        raise OperationError("No document open.")


@dataclass
class PdfToDocxOperation(Operation):
    """Converts the working PDF to a Word document via a text-only
    reconstruction (pdfplumber page text -> python-docx paragraphs,
    one page break per source page). Pure-Python only, no LibreOffice
    attempt - confirmed by hand against the real `soffice` binary that
    a PDF always imports as a *Draw* document, and Draw's export
    filter set (odg, pdf, image formats) does not include docx at all
    ("Error: no export filter for ...sample.docx found, aborting.").
    There is no LibreOffice path here to fall back *from*. The output
    carries no layout/font/image fidelity - it's a text dump, not a
    reconstruction of the original design.
    """

    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_working_pdf(doc)
        assert doc.working_path is not None
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".docx")
        document = docx.Document()
        for i, text in enumerate(extract_pdf_text_by_page(doc.working_path)):
            if i > 0:
                document.add_page_break()  # type: ignore[no-untyped-call]
            for line in text.splitlines() or [""]:
                document.add_paragraph(line)
        document.save(str(out_path))

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {"schema_version": self.schema_version, "type": "pdf_to_docx"}

    def describe(self) -> str:
        return "Converted to Word (text extraction)"


@dataclass
class PdfToPptxOperation(Operation):
    """Converts the working PDF to a PowerPoint presentation by
    rendering each page as a full-slide image at `dpi` (via fitz).
    Pure-Python only, no LibreOffice attempt - confirmed by hand
    against the real `soffice` binary that a PDF always imports as a
    *Draw* document (not Impress), and Draw's export filter set does
    not include pptx either (same "no export filter" failure as
    PdfToDocxOperation - see that class's docstring). There is no
    reliable pure-Python way to reconstruct arbitrary PDF content into
    editable slide shapes either, so this deliberately trades
    editability for visual fidelity instead of guessing at a lossy
    shape reconstruction.
    """

    dpi: int = 150
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def __post_init__(self) -> None:
        if self.dpi < 36:
            raise OperationError(f"dpi must be at least 36, got {self.dpi}.")

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_working_pdf(doc)
        assert doc.working_path is not None
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".pptx")
        with fitz.open(doc.working_path) as src:
            total = src.page_count
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            for i in range(total):
                page_rect = src[i].rect
                prs.slide_width = Emu(round(page_rect.width * _EMU_PER_POINT))
                prs.slide_height = Emu(round(page_rect.height * _EMU_PER_POINT))
                image_path = doc.working_path.parent / f"pdf_to_pptx_p{i}.png"
                render_pdf_page_to_image(doc.working_path, i, self.dpi, image_path)
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image_path),
                    0,
                    0,
                    width=Emu(round(page_rect.width * _EMU_PER_POINT)),
                    height=Emu(round(page_rect.height * _EMU_PER_POINT)),
                )
        prs.save(str(out_path))

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {"schema_version": self.schema_version, "type": "pdf_to_pptx", "dpi": self.dpi}

    def describe(self) -> str:
        return "Converted to PowerPoint (page images)"


@dataclass
class PdfToXlsxOperation(Operation):
    """Extracts every detected table from the working PDF into an
    Excel workbook, one worksheet per page's tables (pdfplumber ->
    openpyxl). Pure-Python only, no LibreOffice attempt - Calc has no
    PDF-import filter at all, so there is nothing to try there first.
    Tabular content only; non-tabular page content is not captured.
    """

    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_working_pdf(doc)
        assert doc.working_path is not None
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".xlsx")
        workbook = openpyxl.Workbook()
        default_sheet = workbook.active
        # openpyxl ships no first-party type stubs for Workbook.active
        # (it's inferred from the untyped source as the generic
        # `_WorkbookChild` base, which has no `.append()`) - the
        # runtime value for a freshly-created Workbook is always a
        # real Worksheet, confirmed by hand; this isinstance check
        # makes that a real runtime guarantee, not just a type: ignore,
        # and narrows for mypy at the same time. Same class of
        # third-party-stub gap CLAUDE.md already documents for
        # scikit-image/tifffile in Phase 4.
        assert isinstance(default_sheet, Worksheet)
        wrote_any = False
        for page_index, tables in enumerate(extract_pdf_tables_by_page(doc.working_path), start=1):
            for table_index, table in enumerate(tables, start=1):
                title = f"Page{page_index}_T{table_index}"[:31]
                sheet = default_sheet if not wrote_any else workbook.create_sheet(title=title)
                if not wrote_any:
                    sheet.title = title
                for row in table:
                    sheet.append(["" if cell is None else cell for cell in row])
                wrote_any = True
        workbook.save(out_path)

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {"schema_version": self.schema_version, "type": "pdf_to_xlsx"}

    def describe(self) -> str:
        return "Converted to Excel (table extraction)"


@dataclass
class PdfToHtmlOperation(Operation):
    """Exports the working PDF's text as a simple HTML file, one
    `<div class="page">` per source page. Pure-Python only, no
    LibreOffice attempt - LibreOffice's PDF -> HTML export rasterizes
    each page as an image inside an HTML wrapper rather than producing
    real markup, so there's nothing useful to try there either. Text
    only; no layout, images, or vector content preserved.
    """

    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_working_pdf(doc)
        assert doc.working_path is not None
        self._pre_snapshot = read_working_bytes(doc)

        out_path = allocate_working_path(doc, suffix=".html")
        parts = ["<!doctype html>", '<html><head><meta charset="utf-8"></head><body>']
        for text in extract_pdf_text_by_page(doc.working_path):
            parts.append('<div class="page">')
            for line in text.splitlines():
                parts.append(f"<p>{html.escape(line)}</p>")
            parts.append("</div>")
        parts.append("</body></html>")
        out_path.write_text("\n".join(parts), encoding="utf-8")

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {"schema_version": self.schema_version, "type": "pdf_to_html"}

    def describe(self) -> str:
        return "Converted to HTML (text extraction)"


@dataclass
class PdfToJpgOperation(Operation):
    """Renders a single page (1-indexed) of the working PDF to a JPEG
    at `dpi`. Deliberately single-page, single-file-out - mirrors
    SignOperation's explicit-page shape and this codebase's existing
    precedent (merge_split.py's module docstring) that producing many
    output files from one document is a batch/export concern layered
    on top in the CLI/GUI, not something a single Operation does.
    """

    page: int
    dpi: int = 200
    _pre_snapshot: bytes | None = field(default=None, init=False, repr=False)

    def __post_init__(self) -> None:
        if self.dpi < 36:
            raise OperationError(f"dpi must be at least 36, got {self.dpi}.")

    def apply(self, doc: DocumentSession) -> DocumentSession:
        _require_working_pdf(doc)
        assert doc.working_path is not None
        self._pre_snapshot = read_working_bytes(doc)

        with fitz.open(doc.working_path) as src:
            total = src.page_count
        if not (1 <= self.page <= total):
            raise OperationError(f"Page {self.page} is out of range (document has {total} pages).")

        out_path = allocate_working_path(doc, suffix=".jpg")
        render_pdf_page_to_image(doc.working_path, self.page - 1, self.dpi, out_path)

        return next_session(doc, out_path)

    def invert(self) -> Operation:
        return snapshot_restore_invert(self._pre_snapshot, self.describe())

    def serialize(self) -> dict[str, Any]:
        return {
            "schema_version": self.schema_version,
            "type": "pdf_to_jpg",
            "page": self.page,
            "dpi": self.dpi,
        }

    def describe(self) -> str:
        return f"Exported page {self.page} to JPEG"


class PdfToDocxPlugin(ToolPlugin):
    tool_id = "pdf_to_docx"
    display_name = "Convert to Word"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        return PdfToDocxOperation()

    def operation_class(self) -> type[Operation]:
        return PdfToDocxOperation


class PdfToPptxPlugin(ToolPlugin):
    tool_id = "pdf_to_pptx"
    display_name = "Convert to PowerPoint"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        return PdfToPptxOperation(dpi=kwargs.get("dpi", 150))

    def operation_class(self) -> type[Operation]:
        return PdfToPptxOperation


class PdfToXlsxPlugin(ToolPlugin):
    tool_id = "pdf_to_xlsx"
    display_name = "Convert to Excel"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        return PdfToXlsxOperation()

    def operation_class(self) -> type[Operation]:
        return PdfToXlsxOperation


class PdfToHtmlPlugin(ToolPlugin):
    tool_id = "pdf_to_html"
    display_name = "Convert to HTML"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        return PdfToHtmlOperation()

    def operation_class(self) -> type[Operation]:
        return PdfToHtmlOperation


class PdfToJpgPlugin(ToolPlugin):
    tool_id = "pdf_to_jpg"
    display_name = "Convert to JPG"
    compatible_core_version = CORE_VERSION_RANGE

    def build_operation(self, **kwargs: Any) -> Operation:
        try:
            page = kwargs["page"]
        except KeyError as exc:
            raise OperationError("Convert to JPG requires a 'page'.") from exc
        return PdfToJpgOperation(page=page, dpi=kwargs.get("dpi", 200))

    def operation_class(self) -> type[Operation]:
        return PdfToJpgOperation
